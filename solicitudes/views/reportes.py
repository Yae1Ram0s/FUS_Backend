import calendar
import os
import uuid
from collections import Counter, defaultdict
from datetime import date, datetime, time, timedelta
from io import BytesIO

from django.conf import settings
from django.db.models import Q
from django.http import FileResponse, Http404, HttpResponse
from django.shortcuts import get_object_or_404
from django.utils import timezone
from django.utils.dateparse import parse_date
from rest_framework.exceptions import PermissionDenied
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from rest_framework.views import APIView

from autenticacion.models import CorreoAutorizado
from ..models import FUS, ReporteGuardado
from ..utils import get_rol, resolver_nombre, _propietario_fus

# Misma paleta que el frontend (Reportes.jsx `COLORES`), para que las
# gráficas de los exportados se vean consistentes con lo que se ve en pantalla.
PALETA_HEX = ['1F5647', 'C9A227', '4DA3E0', 'A35FB0', 'E0805F', '5FAE3F', '8A93A8']
VERDE_INSTITUCIONAL = '1F5647'  # sin '#' — así sirve tal cual para openpyxl (PatternFill) y con '#' antepuesto para reportlab/pptx
LOGO_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'assets', 'logo_hacienda_aduanas.png')


SECCIONES = {
    'resumen': 'Resumen ejecutivo',
    'evolucion': 'Evolución mensual de FUS',
    'estados': 'FUS por estado',
    'carga': 'Carga de trabajo por responsable',
    'unidades': 'FUS por unidad administrativa',
    'productividad': 'Análisis de productividad',
    'tiempos': 'Distribución de tiempos de respuesta',
    'detalle': 'Detalle de solicitudes',
}

MESES_ES = ['ene', 'feb', 'mar', 'abr', 'may', 'jun', 'jul', 'ago', 'sep', 'oct', 'nov', 'dic']

# Etiquetas legibles para los nombres crudos de campo (snake_case) que
# aparecen como encabezado de columna/indicador en los tres exportables —
# un simple `.title()` deja cosas como "Fus Por Dia" o "Cumplimiento Sla",
# que no se ven "nivel ejecutivo". Se comparte entre PDF/Excel/PPTX para que
# los tres digan lo mismo con las mismas palabras.
ETIQUETAS_CAMPO = {
    'total': 'Total de FUS',
    'concluidos': 'Concluidos',
    'pendientes': 'Pendientes',
    'vencidos': 'Vencidos',
    'tiempo_promedio_respuesta': 'Tiempo promedio de respuesta (días)',
    'tiempo_promedio_conclusion': 'Tiempo promedio de conclusión (días)',
    'cumplimiento_sla': 'Cumplimiento SLA (%)',
    'tasa_conclusion': 'Tasa de conclusión (%)',
    'fus_por_dia': 'Productividad (FUS por día)',
    'eficiencia_sla': 'Eficiencia SLA (%)',
    'periodo': 'Periodo',
    'registrados': 'Registrados',
    'nombre': 'Estado',
    'cantidad': 'Cantidad',
    'porcentaje': 'Porcentaje',
    'responsable': 'Responsable',
    'asignados': 'Asignados',
    'capacidad': 'Capacidad (%)',
    'unidad': 'Unidad administrativa',
    'rango': 'Rango de tiempo',
    'folio': 'Folio',
    'fecha_registro': 'Fecha de registro',
    'estado': 'Estado',
    'prioridad': 'Prioridad',
    'fecha_limite': 'Fecha límite',
    'fecha_conclusion': 'Fecha de conclusión',
    'vencido': 'Vencido',
}


def _etiqueta_campo(clave):
    return ETIQUETAS_CAMPO.get(clave, str(clave).replace('_', ' ').title())


def _queryset_usuario(user):
    rol = get_rol(user)
    if rol not in ('ROL1', 'EQUIPO_PARTICULAR'):
        raise PermissionDenied('No autorizado para consultar reportes.')
    qs = FUS.objects.filter(activo=1).select_related(
        'idSolicitanteInterno', 'idComisionado', 'estatusParticular',
        'idMedioRecepcion',
    ).prefetch_related('turnados__idDestinatario')
    if rol == 'EQUIPO_PARTICULAR':
        propietario = _propietario_fus(user)
        return qs.filter(idSolicitanteInterno=propietario) if propietario else qs.none()
    return qs


def _fecha_inicio_fin(params):
    inicio = parse_date(params.get('fecha_inicio', '') or '')
    fin = parse_date(params.get('fecha_fin', '') or '')
    hoy = timezone.localdate()
    if not inicio:
        inicio = hoy.replace(month=1, day=1)
    if not fin:
        fin = hoy
    zona = timezone.get_current_timezone()
    return (
        timezone.make_aware(datetime.combine(inicio, time.min), zona),
        timezone.make_aware(datetime.combine(fin, time.max), zona),
    )


def _aplicar_filtros(qs, params):
    inicio, fin = _fecha_inicio_fin(params)
    qs = qs.filter(fechaRegistro__range=(inicio, fin))
    if params.get('estatus'):
        qs = qs.filter(estatusParticular_id=params['estatus'])
    if params.get('prioridad'):
        qs = qs.filter(prioridad=params['prioridad'])
    if params.get('responsable'):
        responsable = params['responsable']
        qs = qs.filter(
            Q(idComisionado_id=responsable) |
            Q(turnados__idDestinatario_id=responsable, turnados__activo=1)
        ).distinct()
    unidad = params.get('unidad')
    if unidad:
        correos = CorreoAutorizado.objects.filter(
            unidadAdministrativa_id=unidad, activo=1
        ).values_list('email', flat=True)
        qs = qs.filter(idSolicitanteInterno__email__in=correos)
    return qs.order_by('fechaRegistro'), inicio, fin


def _unidad_por_correo():
    return {
        item['email']: {
            'id': item['unidadAdministrativa_id'],
            'nombre': item['unidadAdministrativa__unidadAdministrativa'] or 'Sin unidad',
        }
        for item in CorreoAutorizado.objects.filter(activo=1).values(
            'email', 'unidadAdministrativa_id',
            'unidadAdministrativa__unidadAdministrativa',
        )
    }


def _responsable(fus):
    if fus.idComisionado:
        return fus.idComisionado
    turnados = [t for t in fus.turnados.all() if t.activo and t.idDestinatario]
    return max(turnados, key=lambda t: t.fechaHoraTurnado or t.fechaRegistro).idDestinatario if turnados else None


def _construir_reporte(qs, inicio, fin):
    registros = list(qs)
    ahora = timezone.now()
    estados = Counter()
    meses = defaultdict(lambda: {'registrados': 0, 'concluidos': 0, 'vencidos': 0})
    responsables = defaultdict(lambda: {'asignados': 0, 'pendientes': 0, 'vencidos': 0, 'concluidos': 0})
    unidades = Counter()
    tiempos = Counter({'hasta_1': 0, '1_3': 0, '3_7': 0, 'mas_7': 0})
    mapa_unidades = _unidad_por_correo()
    total_dias = 0
    concluidos_con_tiempo = 0
    total_dias_respuesta = 0
    con_respuesta = 0
    detalle = []

    for fus in registros:
        estado = fus.estatusParticular_id or 'Sin estado'
        estados[estado] += 1
        clave_mes = timezone.localtime(fus.fechaRegistro).strftime('%Y-%m') if fus.fechaRegistro else 'Sin fecha'
        meses[clave_mes]['registrados'] += 1
        vencido = bool(fus.fechaLimite and fus.fechaLimite < ahora and estado != 'Concluido')
        if vencido:
            meses[clave_mes]['vencidos'] += 1
        if fus.fechaConclusion:
            meses[timezone.localtime(fus.fechaConclusion).strftime('%Y-%m')]['concluidos'] += 1
            if fus.fechaRegistro:
                dias = max((fus.fechaConclusion - fus.fechaRegistro).total_seconds() / 86400, 0)
                total_dias += dias
                concluidos_con_tiempo += 1
                if dias <= 1:
                    tiempos['hasta_1'] += 1
                elif dias <= 3:
                    tiempos['1_3'] += 1
                elif dias <= 7:
                    tiempos['3_7'] += 1
                else:
                    tiempos['mas_7'] += 1

        # "Tiempo de respuesta" real = del registro del FUS al primer turnado
        # (primera acción sobre la solicitud), no un valor inventado — si
        # nunca se turnó, no cuenta para el promedio.
        turnados_fus = [t for t in fus.turnados.all() if t.activo]
        if turnados_fus and fus.fechaRegistro:
            primer_turnado = min(turnados_fus, key=lambda t: t.fechaHoraTurnado or t.fechaRegistro)
            fecha_primer_turnado = primer_turnado.fechaHoraTurnado or primer_turnado.fechaRegistro
            if fecha_primer_turnado:
                dias_respuesta = max((fecha_primer_turnado - fus.fechaRegistro).total_seconds() / 86400, 0)
                total_dias_respuesta += dias_respuesta
                con_respuesta += 1

        responsable = _responsable(fus)
        nombre_responsable = resolver_nombre(responsable) if responsable else 'Sin responsable'
        carga = responsables[nombre_responsable]
        carga['asignados'] += 1
        carga['concluidos' if estado == 'Concluido' else 'pendientes'] += 1
        carga['vencidos'] += int(vencido)

        # Unidad de quien atiende el FUS (comisionado o titular turnado —
        # mismo `responsable` de la carga de trabajo), no la del Particular
        # que lo registra: una sola oficina administra solicitudes de
        # distintas direcciones y aduanas, así que su propia unidad no dice
        # nada útil aquí — la de quien realmente lo atiende, sí.
        unidad = mapa_unidades.get(
            responsable.email if responsable else '',
            {'nombre': 'Sin unidad'},
        )['nombre']
        unidades[unidad] += 1
        detalle.append({
            'folio': fus.folio,
            'fecha_registro': timezone.localtime(fus.fechaRegistro).strftime('%d/%m/%Y') if fus.fechaRegistro else '',
            'estado': estado,
            'prioridad': fus.prioridad or 'Sin prioridad',
            'unidad': unidad,
            'responsable': nombre_responsable,
            'fecha_limite': timezone.localtime(fus.fechaLimite).strftime('%d/%m/%Y') if fus.fechaLimite else '',
            'fecha_conclusion': timezone.localtime(fus.fechaConclusion).strftime('%d/%m/%Y') if fus.fechaConclusion else '',
            'vencido': vencido,
        })

    total = len(registros)
    concluidos = estados.get('Concluido', 0)
    pendientes = total - concluidos
    vencidos = sum(1 for item in detalle if item['vencido'])
    dentro_sla = sum(
        1 for fus in registros
        if fus.fechaConclusion and (not fus.fechaLimite or fus.fechaConclusion <= fus.fechaLimite)
    )
    return {
        'periodo': {'inicio': inicio.date().isoformat(), 'fin': fin.date().isoformat()},
        'resumen': {
            'total': total, 'concluidos': concluidos, 'pendientes': pendientes,
            'vencidos': vencidos,
            'tiempo_promedio_respuesta': round(total_dias_respuesta / con_respuesta, 1) if con_respuesta else 0,
            'tiempo_promedio_conclusion': round(total_dias / concluidos_con_tiempo, 1) if concluidos_con_tiempo else 0,
            'cumplimiento_sla': round(dentro_sla * 100 / concluidos, 1) if concluidos else 0,
            'tasa_conclusion': round(concluidos * 100 / total, 1) if total else 0,
        },
        'evolucion': [{'periodo': k, **v} for k, v in sorted(meses.items())],
        'estados': [{'nombre': k, 'cantidad': v, 'porcentaje': round(v * 100 / total, 1) if total else 0} for k, v in estados.most_common()],
        'carga': [{'responsable': k, **v, 'capacidad': min(round(v['pendientes'] * 100 / 20), 100)} for k, v in sorted(responsables.items(), key=lambda x: x[1]['pendientes'], reverse=True)],
        'unidades': [{'unidad': k, 'cantidad': v, 'porcentaje': round(v * 100 / total, 1) if total else 0} for k, v in unidades.most_common()],
        'productividad': {
            'tasa_conclusion': round(concluidos * 100 / total, 1) if total else 0,
            'fus_por_dia': round(total / max((fin.date() - inicio.date()).days + 1, 1), 1),
            'eficiencia_sla': round(dentro_sla * 100 / concluidos, 1) if concluidos else 0,
        },
        'tiempos': [
            {'rango': '≤ 1 día', 'cantidad': tiempos['hasta_1']},
            {'rango': '1 - 3 días', 'cantidad': tiempos['1_3']},
            {'rango': '3 - 7 días', 'cantidad': tiempos['3_7']},
            {'rango': '> 7 días', 'cantidad': tiempos['mas_7']},
        ],
        'detalle': detalle,
    }


# ── Comparación contra un periodo anterior ("vs Abr 2026") ──────────────────

def _restar_meses(anio, mes, n):
    """(año, mes) - n meses, con acarreo de año."""
    indice = (anio * 12 + (mes - 1)) - n
    return indice // 12, indice % 12 + 1


def _rango_mes(anio, mes):
    return date(anio, mes, 1), date(anio, mes, calendar.monthrange(anio, mes)[1])


def _periodo_anterior(inicio, fin, comparar_con=None):
    """Si `comparar_con` trae un mes explícito ("YYYY-MM", lo que manda el
    selector "Comparar con" del frontend), se usa ese mes completo. Si no:
    cuando [inicio, fin] es un bloque de N meses de calendario completos
    (Mes/Trimestre/Semestre/Año ya vienen así desde el frontend), el periodo
    anterior natural son los N meses de calendario inmediatamente antes,
    alineados (ej. mayo 2026 → "vs abril 2026"; Q2 2026 → "vs Q1 2026") — no
    un rango de igual duración pero desalineado del calendario. Para rangos
    personalizados que no calzan con meses completos, sí se usa la misma
    duración inmediatamente antes, que es lo único que tiene sentido ahí."""
    zona = timezone.get_current_timezone()
    if comparar_con:
        try:
            anio, mes = (int(x) for x in comparar_con.split('-'))
            ini_mes, fin_mes = _rango_mes(anio, mes)
            return (
                timezone.make_aware(datetime.combine(ini_mes, time.min), zona),
                timezone.make_aware(datetime.combine(fin_mes, time.max), zona),
            )
        except (ValueError, TypeError):
            pass

    i, f = inicio.date(), fin.date()
    ultimo_dia_f = calendar.monthrange(f.year, f.month)[1]
    if i.day == 1 and f.day == ultimo_dia_f:
        n_meses = (f.year - i.year) * 12 + (f.month - i.month) + 1
        anio_fin_prev, mes_fin_prev = _restar_meses(i.year, i.month, 1)
        anio_ini_prev, mes_ini_prev = _restar_meses(i.year, i.month, n_meses)
        ini_mes_prev, _ = _rango_mes(anio_ini_prev, mes_ini_prev)
        _, fin_mes_prev = _rango_mes(anio_fin_prev, mes_fin_prev)
        return (
            timezone.make_aware(datetime.combine(ini_mes_prev, time.min), zona),
            timezone.make_aware(datetime.combine(fin_mes_prev, time.max), zona),
        )

    duracion = fin - inicio
    fin_prev = inicio - timedelta(microseconds=1)
    inicio_prev = fin_prev - duracion
    return inicio_prev, fin_prev


def _etiqueta_periodo(inicio, fin):
    i, f = inicio.date(), fin.date()
    ultimo_dia_mes = calendar.monthrange(f.year, f.month)[1]
    es_mes_completo = i.day == 1 and f.day == ultimo_dia_mes and i.month == f.month and i.year == f.year
    if es_mes_completo:
        return f"{MESES_ES[i.month - 1].capitalize()} {i.year}"
    return f"{i.strftime('%d/%m')} – {f.strftime('%d/%m/%Y')}"


def _delta_pct(actual, previo):
    if previo:
        return round((actual - previo) * 100 / previo, 1)
    return 0.0 if actual == 0 else 100.0


CAMPOS_COMPARABLES = [
    'total', 'concluidos', 'pendientes', 'vencidos',
    'tiempo_promedio_respuesta', 'tiempo_promedio_conclusion',
    'cumplimiento_sla', 'tasa_conclusion',
]


def _comparar_resumen(actual, previo):
    return {campo: _delta_pct(actual[campo], previo[campo]) for campo in CAMPOS_COMPARABLES}


def _datos(request):
    qs_base = _queryset_usuario(request.user)
    qs, inicio, fin = _aplicar_filtros(qs_base, request.query_params)
    reporte = _construir_reporte(qs, inicio, fin)

    inicio_prev, fin_prev = _periodo_anterior(inicio, fin, request.query_params.get('comparar_con'))
    filtros_prev = request.query_params.dict()
    filtros_prev['fecha_inicio'] = inicio_prev.date().isoformat()
    filtros_prev['fecha_fin'] = fin_prev.date().isoformat()
    qs_prev, _, _ = _aplicar_filtros(qs_base, filtros_prev)
    reporte_prev = _construir_reporte(qs_prev, inicio_prev, fin_prev)

    reporte['comparacion'] = {
        'etiqueta': _etiqueta_periodo(inicio_prev, fin_prev),
        # deltas: variación porcentual por métrica. anterior: valores crudos
        # del periodo previo, para que el frontend pueda mostrar diferencias
        # absolutas donde tenga más sentido (ej. "-1.2 días" en vez de "-18%").
        'deltas': _comparar_resumen(reporte['resumen'], reporte_prev['resumen']),
        'anterior': {campo: reporte_prev['resumen'][campo] for campo in CAMPOS_COMPARABLES},
    }
    return reporte


class ReporteResumenView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        return Response(_datos(request))


class ReporteOpcionesView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        _queryset_usuario(request.user)
        unidades = CorreoAutorizado.objects.filter(
            activo=1, unidadAdministrativa__isnull=False
        ).values('unidadAdministrativa_id', 'unidadAdministrativa__unidadAdministrativa').distinct()
        responsables = CorreoAutorizado.objects.filter(
            activo=1, rol__in=('ROL2', 'COMISIONADO')
        ).order_by('nombre').values('email', 'nombre')
        from django.contrib.auth.models import User
        ids = {u.email: u.id for u in User.objects.filter(email__in=[r['email'] for r in responsables])}
        return Response({
            'secciones': [{'id': k, 'nombre': v} for k, v in SECCIONES.items()],
            'unidades': [{'id': x['unidadAdministrativa_id'], 'nombre': x['unidadAdministrativa__unidadAdministrativa']} for x in unidades],
            'responsables': [{'id': ids.get(x['email']), 'nombre': x['nombre']} for x in responsables if ids.get(x['email'])],
        })


def _secciones(request):
    solicitadas = request.data.get('secciones') or list(SECCIONES)
    return [s for s in solicitadas if s in SECCIONES]


# ── Guardar un reporte exportado (snapshot en disco + registro en BD) ───────

EXTENSIONES = {'pdf': 'pdf', 'excel': 'xlsx', 'pptx': 'pptx'}


def _guardar_si_corresponde(request, user, formato, contenido_bytes, secciones):
    """Si el body trae `guardar: true`, escribe el archivo ya generado a
    MEDIA_ROOT y crea el registro en ReporteGuardado — el snapshot exacto
    que se descargará después no vuelve a calcularse con datos más nuevos."""
    if not request.data.get('guardar'):
        return
    nombre = (request.data.get('nombre') or 'Reporte de solicitudes FUS').strip()[:200]
    extension = EXTENSIONES[formato]
    nombre_fisico = f"{uuid.uuid4().hex}.{extension}"
    ruta_rel = f"reportes/{user.id}/{nombre_fisico}"
    ruta_abs = os.path.join(settings.MEDIA_ROOT, ruta_rel)
    os.makedirs(os.path.dirname(ruta_abs), exist_ok=True)
    with open(ruta_abs, 'wb') as destino:
        destino.write(contenido_bytes)
    ReporteGuardado.objects.create(
        idUsuario=user,
        nombre=nombre,
        formato=formato,
        filtros={k: v for k, v in request.query_params.items()},
        secciones=secciones,
        nombreArchivo=f"{nombre}.{extension}",
        rutaArchivo=ruta_rel,
    )


# ── Gráficas nativas en el Excel exportado (mismas 4 secciones que se grafican
# en pantalla: evolución, estados, unidades, tiempos — carga/detalle son
# tablas también en el frontend, no llevan gráfica ahí tampoco). ──────────────

def _excel_agregar_grafica(ws, clave, n_filas):
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference

    ultima_fila = n_filas + 1  # +1 por el encabezado en la fila 1
    ancla = f"A{ultima_fila + 3}"

    if clave == 'evolucion':
        chart = LineChart()
        chart.title = 'Evolución mensual de FUS'
        chart.style = 2
        chart.height, chart.width = 9, 18
        datos_ref = Reference(ws, min_col=2, max_col=4, min_row=1, max_row=ultima_fila)
        chart.add_data(datos_ref, titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ultima_fila))
        for serie, color in zip(chart.series, PALETA_HEX):
            serie.graphicalProperties.line.solidFill = color
            serie.graphicalProperties.line.width = 22000
        ws.add_chart(chart, ancla)
    elif clave == 'estados':
        from openpyxl.chart.marker import DataPoint
        from openpyxl.chart.shapes import GraphicalProperties
        chart = PieChart()
        chart.title = 'FUS por estado'
        chart.height, chart.width = 9, 14
        chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=ultima_fila), titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ultima_fila))
        chart.series[0].data_points = [
            DataPoint(idx=i, spPr=GraphicalProperties(solidFill=PALETA_HEX[i % len(PALETA_HEX)]))
            for i in range(n_filas)
        ]
        ws.add_chart(chart, ancla)
    elif clave == 'unidades':
        chart = BarChart()
        chart.title = 'FUS por unidad administrativa'
        chart.type = 'bar'
        chart.style = 10
        chart.height, chart.width = 9, 18
        chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=ultima_fila), titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ultima_fila))
        chart.series[0].graphicalProperties.solidFill = PALETA_HEX[0]
        ws.add_chart(chart, ancla)
    elif clave == 'tiempos':
        chart = BarChart()
        chart.title = 'Distribución de tiempos de respuesta'
        chart.style = 10
        chart.height, chart.width = 9, 14
        chart.add_data(Reference(ws, min_col=2, min_row=1, max_row=ultima_fila), titles_from_data=True)
        chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ultima_fila))
        chart.series[0].graphicalProperties.solidFill = PALETA_HEX[1]
        ws.add_chart(chart, ancla)


# Columnas cuyo valor numérico ya está expresado como "42.9" queriendo decir
# "42.9%" (no como fracción 0.429) — se les pone un formato de número que solo
# AGREGA el símbolo "%" como texto, sin que Excel reinterprete/multiplique el
# valor (el formato de porcentaje nativo de Excel sí multiplica x100, lo que
# aquí duplicaría el punto decimal: "42.9" se vería como "4290.0%").
_CAMPOS_PORCENTAJE = {'porcentaje', 'capacidad', 'cumplimiento_sla', 'tasa_conclusion', 'eficiencia_sla'}
_FORMATO_PORCENTAJE = '0.0"%"'


def _excel_portada(wb, datos, seleccion, request):
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Alignment, Font

    ws = wb.create_sheet('Portada', 0)
    ws.sheet_view.showGridLines = False
    ws.column_dimensions['A'].width = 3
    ws.column_dimensions['B'].width = 90

    if os.path.exists(LOGO_PATH):
        img = XLImage(LOGO_PATH)
        img.width, img.height = 267, 45  # mismo lockup y proporción que el PDF/PPTX, a escala de celda
        ws.add_image(img, 'B2')

    ws['B7'] = 'Reporte de solicitudes FUS'
    ws['B7'].font = Font(bold=True, size=20, color=VERDE_INSTITUCIONAL)
    ws['B8'] = 'Inteligencia operativa — Sistema de Control de Solicitudes'
    ws['B8'].font = Font(size=11, color='666666')

    fila = 10
    filas_info = [('Periodo', f"{datos['periodo']['inicio']} a {datos['periodo']['fin']}")]
    filtros_aplicados = [
        f"{etiqueta}: {request.query_params.get(clave)}"
        for clave, etiqueta in [('estatus', 'Estado'), ('prioridad', 'Prioridad'), ('unidad', 'Unidad'), ('responsable', 'Responsable')]
        if request.query_params.get(clave)
    ]
    if filtros_aplicados:
        filas_info.append(('Filtros aplicados', ' · '.join(filtros_aplicados)))
    filas_info.append(('Secciones incluidas', ', '.join(SECCIONES[s] for s in seleccion)))
    filas_info.append(('Generado el', timezone.localtime(timezone.now()).strftime('%d/%m/%Y %H:%M')))

    for etiqueta, valor in filas_info:
        ws.cell(row=fila, column=2, value=etiqueta).font = Font(bold=True, size=10, color=VERDE_INSTITUCIONAL)
        fila += 1
        celda = ws.cell(row=fila, column=2, value=valor)
        celda.font = Font(size=10, color='333333')
        celda.alignment = Alignment(wrap_text=True, vertical='top')
        fila += 2
    return ws


class ReporteExportarExcelView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        import openpyxl
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
        datos = _datos(request)
        seleccion = _secciones(request)
        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        _excel_portada(wb, datos, seleccion, request)

        if 'resumen' in seleccion:
            ws = wb.create_sheet('Resumen')
            ws.append(['Indicador', 'Valor'])
            for clave, valor in datos['resumen'].items():
                ws.append([_etiqueta_campo(clave), valor])
                if clave in _CAMPOS_PORCENTAJE:
                    ws.cell(row=ws.max_row, column=2).number_format = _FORMATO_PORCENTAJE
        tablas = {
            'evolucion': ('Evolución', datos['evolucion']),
            'estados': ('Estados', datos['estados']),
            'carga': ('Carga responsables', datos['carga']),
            'unidades': ('Unidades', datos['unidades']),
            'tiempos': ('Tiempos', datos['tiempos']),
            'detalle': ('Detalle FUS', datos['detalle']),
        }
        for clave, (titulo, filas) in tablas.items():
            if clave not in seleccion:
                continue
            ws = wb.create_sheet(titulo[:31])
            if filas:
                claves_col = list(filas[0])
                ws.append([_etiqueta_campo(k) for k in claves_col])
                for fila in filas:
                    ws.append(list(fila.values()))
                for idx, campo in enumerate(claves_col, start=1):
                    if campo in _CAMPOS_PORCENTAJE:
                        for celda in ws[openpyxl.utils.get_column_letter(idx)][1:]:
                            celda.number_format = _FORMATO_PORCENTAJE
                _excel_agregar_grafica(ws, clave, len(filas))
        if 'productividad' in seleccion:
            ws = wb.create_sheet('Productividad')
            ws.append(['Indicador', 'Valor'])
            for clave, valor in datos['productividad'].items():
                ws.append([_etiqueta_campo(clave), valor])
                if clave in _CAMPOS_PORCENTAJE:
                    ws.cell(row=ws.max_row, column=2).number_format = _FORMATO_PORCENTAJE

        borde = Border(*(Side(style='thin', color='dcdfe0'),) * 4)
        for ws in wb.worksheets:
            if ws.title == 'Portada':
                ws.sheet_properties.tabColor = VERDE_INSTITUCIONAL
                continue
            ws.sheet_properties.tabColor = VERDE_INSTITUCIONAL
            for cell in ws[1]:
                cell.font = Font(bold=True, color='FFFFFF')
                cell.fill = PatternFill('solid', fgColor=VERDE_INSTITUCIONAL)
                cell.alignment = Alignment(vertical='center')
            ws.freeze_panes = 'A2'
            ws.row_dimensions[1].height = 20
            for fila_idx in range(2, ws.max_row + 1):
                if fila_idx % 2 == 0:
                    for celda in ws[fila_idx]:
                        celda.fill = PatternFill('solid', fgColor='F2F6F4')
            for fila in ws.iter_rows(min_row=1, max_row=ws.max_row):
                for celda in fila:
                    celda.border = borde
            for column in ws.columns:
                ws.column_dimensions[column[0].column_letter].width = min(max(len(str(c.value or '')) for c in column) + 3, 45)
        salida = BytesIO()
        wb.save(salida)
        contenido = salida.getvalue()
        _guardar_si_corresponde(request, request.user, 'excel', contenido, seleccion)
        response = HttpResponse(contenido, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        response['Content-Disposition'] = 'attachment; filename="reporte_fus.xlsx"'
        return response


# Misma hoja membretada (logos SHCP/ANAM + pie institucional) que ya usa el
# PDF individual de FUS (ver generar_pdf_fus en views/fus.py) — reusarla es
# lo que hace que este reporte se sienta parte del mismo sistema de
# documentos oficiales, en vez de inventar un membrete propio distinto.
# Es exactamente proporción carta (2550x3300px = 8.5x11in), por eso el PDF
# del reporte va en vertical: en horizontal la imagen saldría deformada.
_LETTERHEAD_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'assets', 'membretada.png')


def _pdf_membrete(canvas_, doc_):
    from reportlab.lib.pagesizes import letter
    canvas_.saveState()
    if os.path.exists(_LETTERHEAD_PATH):
        canvas_.drawImage(_LETTERHEAD_PATH, 0, 0, width=letter[0], height=letter[1], mask='auto', preserveAspectRatio=False)
    canvas_.restoreState()


def _pdf_construir_numbered_canvas():
    """Canvas con doble pasada: dibuja cada página normal en la primera
    pasada, y en showPage() ya conoce el total — así "Página X de Y" es el
    total real del documento, no solo un contador ciego."""
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import cm
    from reportlab.pdfgen import canvas as canvas_mod

    class NumberedCanvas(canvas_mod.Canvas):
        def __init__(self, *args, **kwargs):
            canvas_mod.Canvas.__init__(self, *args, **kwargs)
            self._saved_page_states = []

        def showPage(self):
            self._saved_page_states.append(dict(self.__dict__))
            self._startPage()

        def save(self):
            total = len(self._saved_page_states)
            for state in self._saved_page_states:
                self.__dict__.update(state)
                self._draw_pie_pagina(total)
                canvas_mod.Canvas.showPage(self)
            canvas_mod.Canvas.save(self)

        def _draw_pie_pagina(self, total):
            # y=3.6cm: el bloque "2026 año de Margarita Maza" del membrete
            # arranca ~3.34cm desde abajo (medido sobre membretada.png) — este
            # pie va justo arriba de eso, nunca encima.
            self.saveState()
            self.setFont('Helvetica', 7.5)
            self.setFillColor(_hex_color('#6b6b6b'))
            self.drawRightString(letter[0] - 2 * cm, 3.6 * cm, f"Página {self._pageNumber} de {total}")
            self.drawString(2 * cm, 3.6 * cm, f"Generado el {timezone.localtime(timezone.now()).strftime('%d/%m/%Y %H:%M')} — Sistema de Control de Solicitudes")
            self.restoreState()

    return NumberedCanvas


def _hex_color(hexval):
    from reportlab.lib import colors
    return colors.HexColor(hexval)


def _pdf_grafica(seccion, filas):
    """Dibujo reportlab.graphics equivalente a la gráfica de esa sección en
    pantalla (evolución→línea, estados→pastel, unidades/tiempos→barras) —
    None si la sección no lleva gráfica (ahí solo va la tabla, como en
    pantalla: resumen/carga/productividad/detalle son tarjetas o tablas).
    Ancho acotado a ~450pt: en la hoja carta vertical (con membrete) el
    ancho útil real es ~499pt — pasarse de eso desborda la página. """
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    from reportlab.graphics.charts.legends import Legend
    from reportlab.graphics.charts.linecharts import HorizontalLineChart
    from reportlab.graphics.charts.piecharts import Pie

    if not filas:
        return None
    paleta = [_hex_color(f'#{h}') for h in PALETA_HEX]

    if seccion == 'evolucion':
        d = Drawing(450, 185)
        chart = HorizontalLineChart()
        chart.x, chart.y = 35, 30
        chart.width, chart.height = 300, 135
        chart.data = [
            [f['registrados'] for f in filas],
            [f['concluidos'] for f in filas],
            [f['vencidos'] for f in filas],
        ]
        chart.categoryAxis.categoryNames = [f['periodo'] for f in filas]
        chart.categoryAxis.labels.fontSize = 6.5
        chart.categoryAxis.labels.angle = 30
        chart.valueAxis.valueMin = 0
        # Sin esto, con puros ceros reportlab autoescala a un rango absurdo
        # (0 a 0.01) en vez de mostrar una gráfica plana y legible en 0.
        if max((v for fila in chart.data for v in fila), default=0) == 0:
            chart.valueAxis.valueMax = 1
        etiquetas = ['Registrados', 'Concluidos', 'Vencidos']
        for i in range(3):
            chart.lines[i].strokeColor = paleta[i]
            chart.lines[i].strokeWidth = 2
        d.add(chart)
        leyenda = Legend()
        leyenda.x, leyenda.y = 355, 145
        leyenda.fontSize = 7
        leyenda.colorNamePairs = list(zip(paleta[:3], etiquetas))
        d.add(leyenda)
        return d

    if seccion == 'estados':
        if sum(f['cantidad'] for f in filas) == 0:
            return None
        d = Drawing(450, 185)
        pie = Pie()
        pie.x, pie.y = 160, 15
        pie.width, pie.height = 150, 150
        pie.data = [f['cantidad'] for f in filas]
        pie.labels = [f"{f['nombre']} ({f['porcentaje']}%)" for f in filas]
        pie.slices.strokeWidth = 0.5
        pie.slices.fontSize = 6.5
        for i in range(len(filas)):
            pie.slices[i].fillColor = paleta[i % len(paleta)]
        d.add(pie)
        return d

    if seccion in ('unidades', 'tiempos'):
        clave_label = 'unidad' if seccion == 'unidades' else 'rango'
        d = Drawing(450, 185)
        chart = VerticalBarChart()
        chart.x, chart.y = 40, 40
        chart.width, chart.height = 380, 130
        chart.data = [[f['cantidad'] for f in filas]]
        chart.categoryAxis.categoryNames = [str(f[clave_label])[:16] for f in filas]
        chart.categoryAxis.labels.fontSize = 6.5
        chart.categoryAxis.labels.angle = 20 if seccion == 'unidades' else 0
        chart.categoryAxis.labels.dy = -8
        chart.valueAxis.valueMin = 0
        if max((v for fila in chart.data for v in fila), default=0) == 0:
            chart.valueAxis.valueMax = 1
        chart.bars[0].fillColor = paleta[0 if seccion == 'unidades' else 1]
        d.add(chart)
        return d

    return None


def _pdf_tabla_ejecutiva(tabla_texto, estilo_celda, estilo_celda_header, ancho_total):
    """Tabla con encabezado verde institucional, franjas zebra y celdas
    envueltas en Paragraph (en vez de texto plano) para que el contenido
    largo (nombres de unidad, responsables) haga salto de línea en vez de
    desbordarse o cortarse — la hoja carta vertical deja bastante menos
    ancho útil que la horizontal que usaba antes este reporte."""
    from reportlab.lib import colors
    from reportlab.platypus import Paragraph, Table, TableStyle

    n_cols = len(tabla_texto[0])
    anchos_base = []
    for c in range(n_cols):
        # El encabezado pesa aparte (va en negritas, más ancho por carácter) —
        # medir headers y datos por separado evita que columnas con datos
        # cortos pero encabezado largo (ej. "Prioridad" con valores
        # "Alta"/"Media"/"Baja") terminen partiendo la palabra en 3 líneas.
        largo_header = len(str(tabla_texto[0][c]))
        largo_dato = max((len(str(fila[c])) for fila in tabla_texto[1:]), default=0)
        anchos_base.append(min(max(largo_header * 5.6, largo_dato * 4.6, 48), 175))
    factor = ancho_total / sum(anchos_base)
    col_widths = [a * factor for a in anchos_base]

    filas = [[Paragraph(str(c), estilo_celda_header) for c in tabla_texto[0]]]
    for fila in tabla_texto[1:]:
        filas.append([Paragraph(str(c), estilo_celda) for c in fila])

    t = Table(filas, colWidths=col_widths, repeatRows=1)
    estilos = [
        ('BACKGROUND', (0, 0), (-1, 0), _hex_color('#' + VERDE_INSTITUCIONAL)),
        ('GRID', (0, 0), (-1, -1), .4, colors.HexColor('#d8ded9')),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('TOPPADDING', (0, 0), (-1, -1), 5),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
    ]
    for fila_idx in range(1, len(filas)):
        if fila_idx % 2 == 0:
            estilos.append(('BACKGROUND', (0, fila_idx), (-1, fila_idx), colors.HexColor('#f2f6f4')))
    t.setStyle(TableStyle(estilos))
    return t


class ReporteExportarPDFView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable, KeepTogether

        datos = _datos(request)
        seleccion = _secciones(request)

        st_titulo = ParagraphStyle('rep_titulo', fontName='Helvetica-Bold', fontSize=19, leading=23, textColor=_hex_color('#' + VERDE_INSTITUCIONAL), spaceAfter=8)
        st_subtitulo = ParagraphStyle('rep_subtitulo', fontName='Helvetica', fontSize=9.5, leading=13, textColor=colors.HexColor('#555555'), spaceAfter=2)
        st_seccion = ParagraphStyle('rep_seccion', fontName='Helvetica-Bold', fontSize=12.5, textColor=colors.white, backColor=_hex_color('#' + VERDE_INSTITUCIONAL), leftIndent=0, spaceBefore=0)
        st_celda = ParagraphStyle('rep_celda', fontName='Helvetica', fontSize=7.5, textColor=colors.HexColor('#1a2f28'), leading=9.5)
        st_celda_header = ParagraphStyle('rep_celda_header', parent=st_celda, fontName='Helvetica-Bold', textColor=colors.white)
        st_vacio = ParagraphStyle('rep_vacio', fontName='Helvetica-Oblique', fontSize=8.5, textColor=colors.HexColor('#7a7a7a'))

        ancho_util = letter[0] - 4 * cm  # 2cm de margen a cada lado, igual que el PDF individual de FUS

        filtros_aplicados = [
            f"{etiqueta}: {request.query_params.get(clave)}"
            for clave, etiqueta in [('estatus', 'Estado'), ('prioridad', 'Prioridad'), ('unidad', 'Unidad'), ('responsable', 'Responsable')]
            if request.query_params.get(clave)
        ]

        elementos = [
            Paragraph('Reporte de solicitudes FUS', st_titulo),
            Paragraph('Inteligencia operativa — Sistema de Control de Solicitudes', st_subtitulo),
            Spacer(1, 6),
            HRFlowable(width='100%', thickness=1.4, color=_hex_color('#' + VERDE_INSTITUCIONAL), spaceAfter=8),
            Paragraph(f"<b>Periodo:</b> {datos['periodo']['inicio']} a {datos['periodo']['fin']}", st_subtitulo),
        ]
        if filtros_aplicados:
            elementos.append(Paragraph(f"<b>Filtros aplicados:</b> {' · '.join(filtros_aplicados)}", st_subtitulo))
        elementos.append(Spacer(1, 16))

        contenido = {
            'resumen': datos['resumen'], 'evolucion': datos['evolucion'],
            'estados': datos['estados'], 'carga': datos['carga'],
            'unidades': datos['unidades'], 'productividad': datos['productividad'],
            'tiempos': datos['tiempos'], 'detalle': datos['detalle'][:150],
        }
        for seccion in seleccion:
            filas = contenido[seccion]
            grafica = _pdf_grafica(seccion, filas) if not isinstance(filas, dict) else None

            # Encabezado (+ gráfica, si tiene) siempre juntos: sin esto, un
            # encabezado de sección puede quedar solo hasta abajo de una
            # página con todo su contenido empujado a la siguiente.
            bloque_encabezado = [Paragraph(f"&nbsp;{SECCIONES[seccion]}", st_seccion), Spacer(1, 10)]
            if grafica is not None:
                bloque_encabezado.extend([grafica, Spacer(1, 8)])
            elementos.append(KeepTogether(bloque_encabezado))

            if isinstance(filas, dict):
                tabla = [['Indicador', 'Valor']] + [[_etiqueta_campo(k), v] for k, v in filas.items()]
            elif filas:
                tabla = [[_etiqueta_campo(k) for k in filas[0]]] + [list(x.values()) for x in filas]
            else:
                tabla = None
            if tabla:
                elementos.append(_pdf_tabla_ejecutiva(tabla, st_celda, st_celda_header, ancho_util))
            else:
                elementos.append(Paragraph('Sin información para los filtros seleccionados.', st_vacio))
            elementos.append(Spacer(1, 18))

        salida = BytesIO()
        doc = SimpleDocTemplate(
            salida, pagesize=letter,
            # bottomMargin más alto que el PDF individual de FUS (2.8cm): ahí
            # el membrete no lleva texto encima, pero aquí sí (pie con
            # numeración) y el bloque "2026 año de Margarita Maza" del
            # membrete arranca ~3.34cm desde abajo — con 2.8cm el pie
            # quedaba literalmente escrito encima de ese bloque.
            leftMargin=2 * cm, rightMargin=2 * cm, topMargin=3.3 * cm, bottomMargin=4 * cm,
            title='Reporte de solicitudes FUS', author='Sistema de Control de Solicitudes',
        )
        doc.build(
            elementos,
            onFirstPage=_pdf_membrete, onLaterPages=_pdf_membrete,
            canvasmaker=_pdf_construir_numbered_canvas(),
        )
        contenido_bytes = salida.getvalue()
        _guardar_si_corresponde(request, request.user, 'pdf', contenido_bytes, seleccion)
        response = HttpResponse(contenido_bytes, content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="reporte_fus.pdf"'
        return response


def _pptx_agregar_grafica(slide, seccion, filas):
    """Gráfica nativa de PowerPoint (editable ahí, no una imagen) para las
    mismas 4 secciones que se grafican en pantalla. Devuelve True si agregó
    una gráfica (en ese caso el slide NO lleva también el textbox de datos
    crudos — una gráfica de barras/línea ya comunica la tabla completa)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    if not filas:
        return False

    if seccion == 'evolucion':
        datos_grafica = CategoryChartData()
        datos_grafica.categories = [f['periodo'] for f in filas]
        datos_grafica.add_series('Registrados', [f['registrados'] for f in filas])
        datos_grafica.add_series('Concluidos', [f['concluidos'] for f in filas])
        datos_grafica.add_series('Vencidos', [f['vencidos'] for f in filas])
        slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(1.2), Inches(1.5), Inches(11), Inches(5.3), datos_grafica)
        return True

    if seccion == 'estados':
        datos_grafica = CategoryChartData()
        datos_grafica.categories = [f['nombre'] for f in filas]
        datos_grafica.add_series('FUS', [f['cantidad'] for f in filas])
        slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(3.2), Inches(1.5), Inches(6.6), Inches(5.3), datos_grafica)
        return True

    if seccion in ('unidades', 'tiempos'):
        clave_label = 'unidad' if seccion == 'unidades' else 'rango'
        datos_grafica = CategoryChartData()
        datos_grafica.categories = [f[clave_label] for f in filas]
        datos_grafica.add_series('FUS', [f['cantidad'] for f in filas])
        tipo = XL_CHART_TYPE.BAR_CLUSTERED if seccion == 'unidades' else XL_CHART_TYPE.COLUMN_CLUSTERED
        slide.shapes.add_chart(tipo, Inches(1.2), Inches(1.5), Inches(11), Inches(5.3), datos_grafica)
        return True

    return False


def _pptx_rgb(hexval):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hexval)


def _pptx_banner_titulo(slide, texto, ancho_slide):
    """Mismo lenguaje visual que las secciones del PDF: barra verde
    institucional de borde a borde con el título en blanco."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    barra = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.35), ancho_slide, Inches(0.65))
    barra.fill.solid()
    barra.fill.fore_color.rgb = _pptx_rgb(VERDE_INSTITUCIONAL)
    barra.line.fill.background()
    barra.shadow.inherit = False
    tf = barra.text_frame
    tf.margin_left = Inches(0.35)
    tf.paragraphs[0].text = texto
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = _pptx_rgb('FFFFFF')
    return barra


def _pptx_pie_slide(slide, ancho_slide, alto_slide, numero, total):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    caja = slide.shapes.add_textbox(Inches(0), alto_slide - Inches(0.4), ancho_slide, Inches(0.35))
    p = caja.text_frame.paragraphs[0]
    p.text = f"Sistema de Control de Solicitudes · Página {numero} de {total}"
    p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(9)
    p.runs[0].font.color.rgb = _pptx_rgb('8A8A8A')


def _pptx_agregar_tabla(slide, tabla_texto, x, y, cx, cy):
    """Tabla nativa de PPTX (editable ahí mismo) con el mismo lenguaje visual
    que las tablas del PDF/Excel — reemplaza el textbox de líneas
    "campo: valor" separadas por "|" que traía antes, ilegible en pantalla
    completa y nada "nivel ejecutivo"."""
    from pptx.util import Pt

    filas = tabla_texto[:26]  # una diapositiva solo alcanza para tanto — el resto ya va completo en Excel/PDF
    n_filas, n_cols = len(filas), len(filas[0])
    graphic_frame = slide.shapes.add_table(n_filas, n_cols, x, y, cx, cy)
    tabla = graphic_frame.table
    for c in range(n_cols):
        tabla.columns[c].width = int(cx / n_cols)

    for c, encabezado in enumerate(filas[0]):
        celda = tabla.cell(0, c)
        # texto vacío ("" para fecha_limite/fecha_conclusion sin valor, etc.)
        # no genera ningún run — celda.text_frame.paragraphs[0].runs queda
        # vacío y estilizarlo con runs[0] revienta con IndexError.
        celda.text = str(encabezado) or '—'
        celda.fill.solid()
        celda.fill.fore_color.rgb = _pptx_rgb(VERDE_INSTITUCIONAL)
        p = celda.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].font.bold = True
            p.runs[0].font.size = Pt(12)
            p.runs[0].font.color.rgb = _pptx_rgb('FFFFFF')

    for r in range(1, n_filas):
        for c in range(n_cols):
            celda = tabla.cell(r, c)
            valor = filas[r][c]
            celda.text = str(valor) if valor not in (None, '') else '—'
            celda.fill.solid()
            celda.fill.fore_color.rgb = _pptx_rgb('F2F6F4' if r % 2 == 0 else 'FFFFFF')
            p = celda.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].font.size = Pt(11)
                p.runs[0].font.color.rgb = _pptx_rgb('1A2F28')
    return graphic_frame


class ReporteExportarPPTXView(APIView):
    permission_classes = [IsAuthenticated]

    def post(self, request):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt

        datos = _datos(request)
        seleccion = _secciones(request)
        prs = Presentation()
        ancho, alto = prs.slide_width, prs.slide_height

        # ── Portada de marca ──
        portada = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5])
        for shape in list(portada.shapes):
            if shape.is_placeholder:
                shape.text_frame.clear()

        if os.path.exists(LOGO_PATH):
            from PIL import Image as PILImage
            with PILImage.open(LOGO_PATH) as im:
                proporcion = im.height / im.width
            logo_w = Inches(3.6)
            portada.shapes.add_picture(LOGO_PATH, Inches(0.7), Inches(0.6), width=logo_w, height=int(logo_w * proporcion))

        caja_titulo = portada.shapes.add_textbox(Inches(0.7), Inches(2.6), ancho - Inches(1.4), Inches(1.3))
        p = caja_titulo.text_frame.paragraphs[0]
        p.text = 'Reporte de solicitudes FUS'
        p.runs[0].font.size = Pt(40)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = _pptx_rgb(VERDE_INSTITUCIONAL)

        caja_sub = portada.shapes.add_textbox(Inches(0.7), Inches(3.7), ancho - Inches(1.4), Inches(2))
        tf = caja_sub.text_frame
        tf.word_wrap = True
        lineas_portada = [
            'Inteligencia operativa — Sistema de Control de Solicitudes',
            f"Periodo: {datos['periodo']['inicio']} a {datos['periodo']['fin']}",
        ]
        filtros_aplicados = [
            f"{etiqueta}: {request.query_params.get(clave)}"
            for clave, etiqueta in [('estatus', 'Estado'), ('prioridad', 'Prioridad'), ('unidad', 'Unidad'), ('responsable', 'Responsable')]
            if request.query_params.get(clave)
        ]
        if filtros_aplicados:
            lineas_portada.append('Filtros aplicados: ' + ' · '.join(filtros_aplicados))
        lineas_portada.append(f"Generado el {timezone.localtime(timezone.now()).strftime('%d/%m/%Y %H:%M')}")
        for i, linea in enumerate(lineas_portada):
            par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            par.text = linea
            par.runs[0].font.size = Pt(14)
            par.runs[0].font.color.rgb = _pptx_rgb('595959')
            par.space_after = Pt(6)

        franja = portada.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), alto - Inches(0.35), ancho, Inches(0.35))
        franja.fill.solid()
        franja.fill.fore_color.rgb = _pptx_rgb(VERDE_INSTITUCIONAL)
        franja.line.fill.background()
        franja.shadow.inherit = False

        # ── Slides de contenido ──
        contenido = {
            'resumen': datos['resumen'], 'evolucion': datos['evolucion'],
            'estados': datos['estados'], 'carga': datos['carga'],
            'unidades': datos['unidades'], 'productividad': datos['productividad'],
            'tiempos': datos['tiempos'], 'detalle': datos['detalle'][:25],
        }
        layout_en_blanco = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
        for seccion in seleccion:
            slide = prs.slides.add_slide(layout_en_blanco)
            for shape in list(slide.shapes):
                if shape.is_placeholder:
                    shape.text_frame.clear()
            _pptx_banner_titulo(slide, SECCIONES[seccion], ancho)
            filas = contenido[seccion]

            if not isinstance(filas, dict) and _pptx_agregar_grafica(slide, seccion, filas):
                continue  # el pie con numeración de todas las slides se agrega en un solo pase, después de este for

            if isinstance(filas, dict):
                tabla = [['Indicador', 'Valor']] + [[_etiqueta_campo(k), v] for k, v in filas.items()]
            elif filas:
                tabla = [[_etiqueta_campo(k) for k in filas[0]]] + [list(x.values()) for x in filas]
            else:
                tabla = None

            if tabla:
                _pptx_agregar_tabla(slide, tabla, Inches(0.6), Inches(1.35), ancho - Inches(1.2), alto - Inches(2))
            else:
                caja = slide.shapes.add_textbox(Inches(0.8), Inches(2), ancho - Inches(1.6), Inches(1))
                p = caja.text_frame.paragraphs[0]
                p.text = 'Sin información para los filtros seleccionados.'
                p.runs[0].font.size = Pt(14)
                p.runs[0].font.italic = True
                p.runs[0].font.color.rgb = _pptx_rgb('7A7A7A')

        total_slides = len(prs.slides)
        for i, slide in enumerate(prs.slides, start=1):
            if i == 1:
                continue  # la portada ya lleva su propia franja, sin numeración
            _pptx_pie_slide(slide, ancho, alto, i, total_slides)

        salida = BytesIO()
        prs.save(salida)
        contenido_bytes = salida.getvalue()
        _guardar_si_corresponde(request, request.user, 'pptx', contenido_bytes, seleccion)
        response = HttpResponse(contenido_bytes, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="reporte_fus.pptx"'
        return response


# ── Reportes guardados ───────────────────────────────────────────────────────

class ReporteGuardadoListView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request):
        _queryset_usuario(request.user)
        guardados = ReporteGuardado.objects.filter(idUsuario=request.user)[:50]
        return Response([
            {
                'id': g.id,
                'nombre': g.nombre,
                'formato': g.formato,
                'fechaCreacion': g.fechaCreacion,
                'secciones': g.secciones,
            }
            for g in guardados
        ])


class ReporteGuardadoDescargarView(APIView):
    permission_classes = [IsAuthenticated]

    def get(self, request, pk):
        _queryset_usuario(request.user)
        guardado = get_object_or_404(ReporteGuardado, pk=pk, idUsuario=request.user)

        media_root = os.path.realpath(settings.MEDIA_ROOT)
        ruta = os.path.realpath(os.path.join(media_root, guardado.rutaArchivo))
        if os.path.commonpath([media_root, ruta]) != media_root:
            raise Http404
        if not os.path.exists(ruta):
            raise Http404
        return FileResponse(open(ruta, 'rb'), as_attachment=True, filename=guardado.nombreArchivo)

    def delete(self, request, pk):
        _queryset_usuario(request.user)
        guardado = get_object_or_404(ReporteGuardado, pk=pk, idUsuario=request.user)
        media_root = os.path.realpath(settings.MEDIA_ROOT)
        ruta = os.path.realpath(os.path.join(media_root, guardado.rutaArchivo))
        if os.path.commonpath([media_root, ruta]) == media_root and os.path.exists(ruta):
            os.remove(ruta)
        guardado.delete()
        return Response(status=204)
